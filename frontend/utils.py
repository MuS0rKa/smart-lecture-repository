from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def extract_text(file):
    file_extension = file.name.split('.')[-1].lower()
    text = ""

    if file_extension == 'pdf':
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
            
    elif file_extension == 'docx':
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
            
    elif file_extension == 'pptx':
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                    
    elif file_extension == 'txt':
        text = file.read().decode("utf-8")
        
    return text.strip()